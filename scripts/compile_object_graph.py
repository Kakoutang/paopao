"""Deterministic compiler: object graph JSON -> editable PPTX.

No AI, no heuristics, no DOM extraction. Reads structured JSON,
emits python-pptx shapes. Same JSON -> same PPTX every time.

Usage:
    compile_deck(graph_paths, pptx_path) -> summary dict
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
from lxml.etree import SubElement

FULL_SLIDE_IMAGE_AREA_RATIO = 0.82

# Canvas: 1920x1080 px -> 13.333" x 7.5"
HTML_W = 1920
HTML_H = 1080
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
EMU_PER_IN = 914400
PX_TO_EMU_X = (SLIDE_W_IN * EMU_PER_IN) / HTML_W
PX_TO_EMU_Y = (SLIDE_H_IN * EMU_PER_IN) / HTML_H
PX_TO_PT = (SLIDE_W_IN * 72) / HTML_W

PALETTE = {
    "deep_blue": RGBColor(0x30, 0x54, 0x96),
    "soft_blue": RGBColor(0x5B, 0x9B, 0xD5),
    "consulting_blue": RGBColor(0x44, 0x72, 0xC4),
    "grey": RGBColor(0x66, 0x66, 0x66),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "light_grey": RGBColor(0xF2, 0xF2, 0xF2),
    "border_grey": RGBColor(0xE0, 0xE4, 0xED),
    "text_dark": RGBColor(0x33, 0x33, 0x33),
}

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

CHART_TYPE_MAP = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
}

CHART_ALIASES = {
    "combo": "column",
    "funnel": "bar_stacked",
    "waterfall": "bar_stacked",
    "histogram": "column",
    "donut": "doughnut",
    "stacked_bar": "bar_stacked",
    "stacked_column": "column_stacked",
    "bubble": "scatter",
}

SERIES_COLORS = [
    RGBColor(0x30, 0x54, 0x96),
    RGBColor(0x5B, 0x9B, 0xD5),
    RGBColor(0x44, 0x72, 0xC4),
    RGBColor(0x66, 0x66, 0x66),
]
MAX_SERIES = 4

HEADER_BG = RGBColor(0x28, 0x3E, 0x84)
BORDER_LIGHT = RGBColor(0xE0, 0xE4, 0xED)
BORDER_WIDTH_EMU = 6350  # 0.5pt


def _px(bbox: list) -> tuple[Emu, Emu, Emu, Emu]:
    x, y, w, h = bbox
    return (
        Emu(int(x * PX_TO_EMU_X)),
        Emu(int(y * PX_TO_EMU_Y)),
        Emu(max(1, int(w * PX_TO_EMU_X))),
        Emu(max(1, int(h * PX_TO_EMU_Y))),
    )


def _parse_color(s: str | None) -> RGBColor | None:
    if not s:
        return None
    s = s.strip().lstrip("#")
    if len(s) == 6:
        try:
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
        except ValueError:
            return None
    return None


def _style_color(style: dict, key: str, default: RGBColor | None = None) -> RGBColor | None:
    val = style.get(key)
    if val:
        return _parse_color(val) or default
    return default


def _float_value(value, default: float) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def _strip_theme_style(shp) -> None:
    """Remove <p:style> so shape uses its own spPr, not theme defaults."""
    el = shp._element.find(qn("p:style"))
    if el is not None:
        shp._element.remove(el)


def _rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _pct_to_alpha(value) -> str | None:
    try:
        pct = float(value)
    except (TypeError, ValueError):
        return None
    pct = max(0.0, min(100.0, pct))
    return str(int(round((100.0 - pct) * 1000)))


def _set_solid_fill_alpha(shp, transparency) -> None:
    alpha = _pct_to_alpha(transparency)
    if alpha is None:
        return
    try:
        sp_pr = shp._element.find(qn("p:spPr"))
        solid = sp_pr.find(qn("a:solidFill")) if sp_pr is not None else None
        srgb = solid.find(qn("a:srgbClr")) if solid is not None else None
        if srgb is None:
            return
        for child in list(srgb):
            if child.tag == qn("a:alpha"):
                srgb.remove(child)
        alpha_el = SubElement(srgb, qn("a:alpha"))
        alpha_el.set("val", alpha)
    except Exception:
        pass


def _lock_text_body(tf, style: dict, default_margin: int = 0) -> None:
    tf.word_wrap = bool(style.get("word_wrap", True))
    margin = style.get("padding")
    if isinstance(margin, (int, float)):
        margins = {
            "left": float(margin),
            "right": float(margin),
            "top": float(margin),
            "bottom": float(margin),
        }
    elif isinstance(margin, dict):
        margins = {
            "left": _float_value(margin.get("left"), default_margin),
            "right": _float_value(margin.get("right"), default_margin),
            "top": _float_value(margin.get("top"), default_margin),
            "bottom": _float_value(margin.get("bottom"), default_margin),
        }
    else:
        margins = {
            "left": _float_value(style.get("padding_left"), default_margin),
            "right": _float_value(style.get("padding_right"), default_margin),
            "top": _float_value(style.get("padding_top"), default_margin),
            "bottom": _float_value(style.get("padding_bottom"), default_margin),
        }
    tf.margin_left = Emu(int(margins["left"] * PX_TO_EMU_X))
    tf.margin_right = Emu(int(margins["right"] * PX_TO_EMU_X))
    tf.margin_top = Emu(int(margins["top"] * PX_TO_EMU_Y))
    tf.margin_bottom = Emu(int(margins["bottom"] * PX_TO_EMU_Y))

    bodyPr = tf._txBody.bodyPr
    v_align = style.get("vertical_align", "top")
    if v_align == "center":
        bodyPr.set("anchor", "ctr")
    elif v_align == "bottom":
        bodyPr.set("anchor", "b")
    else:
        bodyPr.set("anchor", "t")
    bodyPr.set("wrap", "square" if tf.word_wrap else "none")

    for child in list(bodyPr):
        tag = child.tag
        if tag.endswith("}noAutofit") or tag.endswith("}spAutoFit") or tag.endswith("}normAutofit"):
            bodyPr.remove(child)
    if style.get("auto_fit") is True:
        SubElement(bodyPr, qn("a:normAutofit"))
    else:
        SubElement(bodyPr, qn("a:noAutofit"))


def _lock_run_typeface(run, font_family: str) -> None:
    try:
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea", "a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = SubElement(rPr, qn(tag))
            el.set("typeface", font_family)
    except Exception:
        pass


def _apply_shadow(shp, style: dict) -> None:
    shadow = style.get("shadow")
    if not shadow:
        return
    if not isinstance(shadow, dict):
        shadow = {}
    color = _parse_color(str(shadow.get("color", style.get("shadow_color", "#000000")))) or RGBColor(0, 0, 0)
    blur = _float_value(shadow.get("blur", style.get("shadow_blur", 6)), 6.0)
    distance = _float_value(shadow.get("distance", style.get("shadow_distance", 2)), 2.0)
    angle = _float_value(shadow.get("angle", style.get("shadow_angle", 45)), 45.0)
    transparency = shadow.get("transparency", style.get("shadow_transparency", 65))
    alpha = _pct_to_alpha(transparency) or "35000"
    try:
        sp_pr = shp._element.find(qn("p:spPr"))
        if sp_pr is None:
            return
        for child in list(sp_pr):
            if child.tag == qn("a:effectLst"):
                sp_pr.remove(child)
        effect = SubElement(sp_pr, qn("a:effectLst"))
        shd = SubElement(effect, qn("a:outerShdw"))
        shd.set("blurRad", str(int(blur * 12700)))
        shd.set("dist", str(int(distance * 12700)))
        shd.set("dir", str(int(angle * 60000)))
        shd.set("algn", "ctr")
        shd.set("rotWithShape", "0")
        clr = SubElement(shd, qn("a:srgbClr"))
        clr.set("val", _rgb_hex(color))
        alpha_el = SubElement(clr, qn("a:alpha"))
        alpha_el.set("val", alpha)
    except Exception:
        pass


def _apply_line_style(line, style: dict, default_color: RGBColor | None = None, default_width_px: float = 1.0) -> None:
    color = _style_color(style, "line_color", _style_color(style, "border_color", default_color))
    if color:
        line.color.rgb = color
    width_px = style.get("line_width", style.get("width", default_width_px))
    line.width = Emu(max(BORDER_WIDTH_EMU, int(_float_value(width_px, default_width_px) * 12700)))
    dash = str(style.get("dash", style.get("line_dash", ""))).strip().lower()
    if dash:
        try:
            ln = line._ln
            for child in list(ln):
                if child.tag == qn("a:prstDash"):
                    ln.remove(child)
            dash_el = SubElement(ln, qn("a:prstDash"))
            dash_el.set("val", {"dashed": "dash", "dash": "dash", "dotted": "dot", "dot": "dot"}.get(dash, dash))
        except Exception:
            pass


def _set_cell_border(tcPr, tag: str, color: RGBColor, width: int = BORDER_WIDTH_EMU) -> None:
    existing = tcPr.find(qn(tag))
    if existing is not None:
        tcPr.remove(existing)
    ln = SubElement(tcPr, qn(tag))
    ln.set("w", str(width))
    sf = SubElement(ln, qn("a:solidFill"))
    clr = SubElement(sf, qn("a:srgbClr"))
    clr.set("val", _rgb_hex(color))


def _set_preset_adjustment(shape, name: str, value: int) -> None:
    try:
        prstGeom = shape._element.find(qn("p:spPr")).find(qn("a:prstGeom"))
        if prstGeom is None:
            return
        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = SubElement(prstGeom, qn("a:avLst"))
        for child in list(avLst):
            if child.tag == qn("a:gd") and child.get("name") == name:
                avLst.remove(child)
        gd = SubElement(avLst, qn("a:gd"))
        gd.set("name", name)
        gd.set("fmla", f"val {int(value)}")
    except Exception:
        pass


def _set_chart_manual_layout(chart, layout: dict) -> None:
    if not isinstance(layout, dict):
        return
    try:
        plot_area = chart._element.find(".//" + qn("c:plotArea"))
        if plot_area is None:
            return
        existing = plot_area.find(qn("c:layout"))
        if existing is not None:
            plot_area.remove(existing)
        layout_el = SubElement(plot_area, qn("c:layout"))
        manual = SubElement(layout_el, qn("c:manualLayout"))
        target = SubElement(manual, qn("c:layoutTarget"))
        target.set("val", "inner")
        for tag in ("xMode", "yMode", "wMode", "hMode"):
            el = SubElement(manual, qn(f"c:{tag}"))
            el.set("val", "factor")
        defaults = {"x": 0.08, "y": 0.14, "w": 0.84, "h": 0.70}
        for tag, default in defaults.items():
            value = _float_value(layout.get(tag), default)
            value = min(1.0, max(0.0, value))
            el = SubElement(manual, qn(f"c:{tag}"))
            el.set("val", f"{value:.4f}")
    except Exception:
        pass


def _set_chart_bar_geometry(chart, gap_width: object = None, overlap: object = None) -> None:
    try:
        plot_area = chart._element.find(".//" + qn("c:plotArea"))
        if plot_area is None:
            return
        for chart_tag in ("barChart", "bar3DChart"):
            bar_chart = plot_area.find(qn(f"c:{chart_tag}"))
            if bar_chart is None:
                continue
            if gap_width is not None:
                gap = _float_value(gap_width, 80)
                gap = int(max(0, min(500, gap)))
                existing = bar_chart.find(qn("c:gapWidth"))
                if existing is not None:
                    bar_chart.remove(existing)
                el = SubElement(bar_chart, qn("c:gapWidth"))
                el.set("val", str(gap))
            if overlap is not None:
                val = int(max(-100, min(100, _float_value(overlap, 0))))
                existing = bar_chart.find(qn("c:overlap"))
                if existing is not None:
                    bar_chart.remove(existing)
                el = SubElement(bar_chart, qn("c:overlap"))
                el.set("val", str(val))
    except Exception:
        pass


def _get_text_lines(obj: dict) -> list[str]:
    text = obj.get("text", "")
    if isinstance(text, list):
        return [str(t) for t in text if str(t).strip()]
    return [line for line in str(text).split("\n") if line.strip()] if text else []


def _get_text_runs(obj: dict) -> list[dict]:
    runs = obj.get("text_runs")
    if not isinstance(runs, list):
        return []
    return [run for run in runs if isinstance(run, dict) and str(run.get("text", "")).strip()]


def _paragraph_line_spacing(value):
    if value is None:
        return None
    try:
        numeric = float(value)
    except (TypeError, ValueError):
        return None
    if 0 < numeric <= 3:
        return numeric
    if numeric > 3:
        return Pt(numeric)
    return None


def _set_paragraph_spacing(paragraph, style: dict) -> None:
    paragraph.space_before = Pt(_float_value(style.get("space_before"), 0.0))
    paragraph.space_after = Pt(_float_value(style.get("space_after"), 0.0))
    spacing = _paragraph_line_spacing(style.get("line_spacing"))
    if spacing is not None:
        paragraph.line_spacing = spacing


def _text_parts(obj: dict) -> tuple[str, list[str], str]:
    component = obj.get("component")
    if isinstance(component, dict):
        metric = str(component.get("metric", component.get("number", ""))).strip()
        desc = component.get("description", component.get("body", []))
        tag = str(component.get("tag", component.get("label", ""))).strip()
        if isinstance(desc, list):
            desc_lines = [str(item).strip() for item in desc if str(item).strip()]
        else:
            desc_lines = [line.strip() for line in str(desc).splitlines() if line.strip()]
        return metric, desc_lines, tag

    lines = _get_text_lines(obj)
    if not lines:
        return "", [], ""
    if len(lines) == 1:
        return lines[0], [], ""
    return lines[0], lines[1:-1], lines[-1]


def _add_text(slide, bbox: list, text, style: dict) -> None:
    lines = text if isinstance(text, list) else [str(text)]
    _add_textbox(slide, bbox, [line for line in lines if str(line).strip()], style)


def _apply_run_style(run, run_data: dict, style: dict) -> None:
    font_family = run_data.get("font_family", style.get("font_family", "Arial"))
    font_family = str(font_family)
    run.font.name = font_family
    _lock_run_typeface(run, font_family)
    color = _parse_color(str(run_data.get("color", ""))) or _style_color(style, "color", PALETTE["text_dark"])
    if color:
        run.font.color.rgb = color
    font_size = run_data.get("font_size", style.get("font_size", 14))
    run.font.size = Pt(_float_value(font_size, 14.0))
    font_weight = str(run_data.get("font_weight", style.get("font_weight", "normal")))
    run.font.bold = bool(run_data.get("bold", False)) or font_weight in ("bold", "700", "800", "900")
    run.font.italic = bool(run_data.get("italic", False)) or str(run_data.get("font_style", "")).lower() == "italic"
    run.font.italic = bool(run_data.get("italic", False))


def _add_textbox(slide, bbox: list, text_lines: list[str], style: dict, text_runs: list[dict] | None = None) -> None:
    text_runs = text_runs or []
    if (not text_lines or all(not t.strip() for t in text_lines)) and not text_runs:
        return
    x, y, w, h = _px(bbox)
    tb = slide.shapes.add_textbox(x, y, w, h)
    _strip_theme_style(tb)
    tf = tb.text_frame
    _lock_text_body(tf, style, default_margin=0)

    font_size = _float_value(style.get("font_size"), 14.0)
    alignment = ALIGN_MAP.get(style.get("text_align", "left"), PP_ALIGN.LEFT)
    line_spacing = style.get("line_spacing")

    if text_runs:
        p = tf.paragraphs[0]
        p.alignment = alignment
        _set_paragraph_spacing(p, style)
        first_run = True
        for run_data in text_runs:
            if not first_run and run_data.get("new_paragraph"):
                p = tf.add_paragraph()
                p.alignment = alignment
                _set_paragraph_spacing(p, style)
            first_run = False
            run = p.add_run()
            run.text = str(run_data.get("text", ""))
            _apply_run_style(run, run_data, style)
            if run_data.get("break_after"):
                p = tf.add_paragraph()
                p.alignment = alignment
                _set_paragraph_spacing(p, style)
        return

    color = _style_color(style, "color", PALETTE["text_dark"])
    font_weight = style.get("font_weight", "normal")
    font_family = style.get("font_family", "Arial")
    first = True
    for line in text_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = alignment
        _set_paragraph_spacing(p, style)
        run = p.add_run()
        run.text = line
        run.font.name = font_family
        _lock_run_typeface(run, str(font_family))
        if color:
            run.font.color.rgb = color
        run.font.size = Pt(font_size)
        run.font.bold = font_weight in ("bold", "700", "800", "900")


def _emit_shape(slide, obj: dict) -> None:
    bbox = obj.get("bbox", [0, 0, 100, 100])
    style = obj.get("style", {})
    x, y, w, h = _px(bbox)
    radius = style.get("border_radius", 0)

    shape_type = str(style.get("shape_type") or obj.get("object_kind") or "").strip().lower()
    if shape_type == "chevron":
        shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
        _set_preset_adjustment(shp, "adj", int(_float_value(style.get("chevron_adjustment"), 42000)))
    elif shape_type == "oval":
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    elif radius > 1:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shorter = min(int(w), int(h)) or 1
        radius_emu = int(radius * PX_TO_EMU_X)
        adj_val = min(50000, int(radius_emu / shorter * 100000))
        try:
            prstGeom = shp._element.find(qn("p:spPr")).find(qn("a:prstGeom"))
            if prstGeom is not None:
                avLst = prstGeom.find(qn("a:avLst"))
                if avLst is None:
                    avLst = SubElement(prstGeom, qn("a:avLst"))
                else:
                    for child in list(avLst):
                        avLst.remove(child)
                gd = SubElement(avLst, qn("a:gd"))
                gd.set("name", "adj")
                gd.set("fmla", f"val {adj_val}")
        except Exception:
            pass
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)

    _strip_theme_style(shp)

    accent = style.get("accent_border")
    accent_only = isinstance(accent, dict) and not _style_color(style, "background") and not style.get("border")

    fill = _style_color(style, "background")
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        _set_solid_fill_alpha(shp, style.get("fill_transparency", style.get("background_transparency")))
    elif accent_only:
        shp.fill.background()
    else:
        shp.fill.background()

    border = style.get("border")
    if isinstance(border, dict):
        bc = _parse_color(border.get("color", ""))
        bw = border.get("width", 1)
        if bc:
            dash = border.get("dash", style.get("dash", ""))
            if not dash and str(style.get("border_style", "")).strip().lower() in {"dash", "dashed", "dotted", "dot"}:
                dash = style.get("border_style")
            _apply_line_style(
                shp.line,
                {**style, "line_color": border.get("color"), "line_width": bw, "dash": dash},
                default_color=bc,
                default_width_px=_float_value(bw, 1.0),
            )
        else:
            shp.line.fill.background()
    elif accent_only:
        shp.line.fill.background()
    else:
        shp.line.fill.background()
    _apply_shadow(shp, style)

    if isinstance(accent, dict):
        _emit_accent_border(slide, bbox, accent)

    text_lines = _get_text_lines(obj)
    text_runs = _get_text_runs(obj)
    if (text_lines and any(t.strip() for t in text_lines)) or text_runs:
        tf = shp.text_frame
        _lock_text_body(tf, {"vertical_align": "center", "padding": 6, **style}, default_margin=6)
        font_size = _float_value(style.get("font_size"), 12.0)
        alignment = ALIGN_MAP.get(style.get("text_align", "center"), PP_ALIGN.CENTER)
        if text_runs:
            p = tf.paragraphs[0]
            p.alignment = alignment
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            first_run = True
            for run_data in text_runs:
                if not first_run and run_data.get("new_paragraph"):
                    p = tf.add_paragraph()
                    p.alignment = alignment
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                first_run = False
                run = p.add_run()
                run.text = str(run_data.get("text", ""))
                _apply_run_style(run, run_data, {"font_size": font_size, **style})
                if run_data.get("break_after"):
                    p = tf.add_paragraph()
                    p.alignment = alignment
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
            return
        color = _style_color(style, "color", PALETTE["text_dark"])
        first = True
        for line in text_lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = alignment
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            run = p.add_run()
            run.text = line
            font_family = str(style.get("font_family", "Arial"))
            run.font.name = font_family
            _lock_run_typeface(run, font_family)
            if color:
                run.font.color.rgb = color
            run.font.size = Pt(font_size)
            run.font.bold = style.get("font_weight", "") in ("bold", "700", "800", "900")
            run.font.italic = str(style.get("font_style", "")).lower() == "italic" or bool(style.get("italic", False))


def _emit_text(slide, obj: dict) -> None:
    bbox = obj.get("bbox", [0, 0, 100, 100])
    style = obj.get("style", {})
    text_lines = _get_text_lines(obj)
    _add_textbox(slide, bbox, text_lines, style, _get_text_runs(obj))


def _emit_kpi(slide, obj: dict) -> None:
    bbox = obj.get("bbox", [0, 0, 320, 180])
    style = obj.get("style", {})
    x, y, w, h = bbox
    metric, desc_lines, tag = _text_parts(obj)

    card_style = {
        "background": style.get("background", "#FFFFFF"),
        "border": style.get("border", {"width": 1, "color": "#5B9BD5"}),
        "border_radius": style.get("border_radius", 5),
        "padding": 0,
    }
    _emit_shape(slide, {**obj, "object_kind": "shape", "text": [], "style": card_style})

    metric_style = {
        "color": style.get("metric_color", style.get("color", "#305496")),
        "font_size": style.get("metric_font_size", max(10, _float_value(style.get("font_size"), 14))),
        "font_weight": style.get("metric_font_weight", "bold"),
        "font_family": style.get("font_family", "Arial"),
        "text_align": "center",
        "vertical_align": "center",
        "padding": 0,
        "line_spacing": style.get("line_spacing", 0.9),
        "auto_fit": False,
    }
    body_style = {
        "color": style.get("body_color", "#1C1917"),
        "font_size": style.get("body_font_size", max(8, _float_value(style.get("font_size"), 12) - 3)),
        "font_weight": style.get("body_font_weight", "normal"),
        "font_family": style.get("font_family", "Arial"),
        "text_align": "center",
        "vertical_align": "center",
        "padding": 0,
        "line_spacing": style.get("body_line_spacing", 0.92),
        "auto_fit": False,
    }
    tag_style = {
        "background": style.get("tag_background", style.get("accent", "#005DAA")),
        "color": style.get("tag_color", "#FFFFFF"),
        "font_size": style.get("tag_font_size", max(8, _float_value(style.get("font_size"), 12) - 3)),
        "font_weight": style.get("tag_font_weight", "normal"),
        "font_family": style.get("font_family", "Arial"),
        "text_align": "center",
        "vertical_align": "center",
        "border": {"width": 0.1, "color": style.get("tag_background", style.get("accent", "#005DAA"))},
        "border_radius": style.get("tag_border_radius", 4),
        "padding": 2,
        "line_spacing": 0.9,
        "auto_fit": False,
    }
    pad_x = _float_value(style.get("inner_padding_x"), 14)
    top = _float_value(style.get("metric_top"), 14)
    metric_h = _float_value(style.get("metric_height"), 32)
    tag_h = _float_value(style.get("tag_height"), 26)
    tag_bottom = _float_value(style.get("tag_bottom"), 12)
    tag_x = x + pad_x
    tag_y = y + h - tag_bottom - tag_h
    _add_text(slide, [x + pad_x, y + top, w - pad_x * 2, metric_h], metric, metric_style)
    _add_text(slide, [x + pad_x, y + top + metric_h + 4, w - pad_x * 2, max(20, tag_y - (y + top + metric_h + 8))], desc_lines, body_style)
    if tag:
        _emit_shape(
            slide,
            {
                **obj,
                "object_kind": "shape",
                "bbox": [tag_x, tag_y, w - pad_x * 2, tag_h],
                "text": tag,
                "style": tag_style,
            },
        )


def _split_callout_text(obj: dict) -> tuple[str, list[str]]:
    component = obj.get("component")
    if isinstance(component, dict):
        number = str(component.get("number", "")).strip()
        body = component.get("body", component.get("text", []))
        body_lines = body if isinstance(body, list) else str(body).splitlines()
        return number, [str(line).strip() for line in body_lines if str(line).strip()]
    lines = _get_text_lines(obj)
    if not lines:
        return "", []
    first = lines[0].strip()
    parts = first.split(maxsplit=1)
    if parts and parts[0].isdigit():
        number = parts[0]
        body_lines = ([parts[1]] if len(parts) > 1 else []) + lines[1:]
        return number, body_lines
    return "", lines


def _split_callout_title(body_lines: list[str]) -> tuple[str, list[str]]:
    if not body_lines:
        return "", []
    first = body_lines[0].strip()
    if ":" not in first:
        return "", body_lines
    title, rest = first.split(":", 1)
    title = title.strip()
    if not title:
        return "", body_lines
    remainder = rest.strip()
    lines = ([remainder] if remainder else []) + body_lines[1:]
    return f"{title}:", [line for line in lines if str(line).strip()]


def _emit_note_icon(slide, bbox: list, style: dict) -> None:
    x, y, _w, h = bbox
    size = min(34, max(18, h - 20))
    icon_x = x + 22
    icon_y = y + (h - size) / 2
    ix, iy, iw, ih = _px([icon_x, icon_y, size, size])
    shp = slide.shapes.add_shape(MSO_SHAPE.CIRCULAR_ARROW, ix, iy, iw, ih)
    _strip_theme_style(shp)
    border = style.get("border")
    border_color = _parse_color(border.get("color", "")) if isinstance(border, dict) else None
    color = _style_color(style, "icon_color", border_color or PALETTE["deep_blue"])
    shp.fill.solid()
    shp.fill.fore_color.rgb = color or PALETTE["deep_blue"]
    shp.line.fill.background()


def _emit_callout(slide, obj: dict) -> None:
    bbox = obj.get("bbox", [0, 0, 320, 120])
    style = obj.get("style", {})
    x, y, w, h = bbox
    number, body_lines = _split_callout_text(obj)
    _emit_shape(
        slide,
        {
            **obj,
            "object_kind": "shape",
            "text": [],
            "style": {
                "background": style.get("background", "#FFFFFF"),
                "border": style.get("border", {"width": 1, "color": "#5B9BD5"}),
                "border_radius": style.get("border_radius", 5),
                "padding": 0,
            },
        },
    )
    badge = style.get("badge", {}) if isinstance(style.get("badge"), dict) else {}
    badge_size = _float_value(badge.get("size", style.get("badge_size")), 34)
    badge_x = x + _float_value(badge.get("x", style.get("badge_x")), 18)
    badge_y = y + _float_value(badge.get("y", style.get("badge_y")), 18)
    is_note_strip = not number and w >= 600 and h <= 90
    if is_note_strip:
        _emit_note_icon(slide, bbox, style)
    if number:
        _emit_shape(
            slide,
            {
                **obj,
                "object_kind": "badge",
                "bbox": [badge_x, badge_y, badge_size, badge_size],
                "text": number,
                "style": {
                    "shape_type": "oval",
                    "background": badge.get("background", "#005DAA"),
                    "color": badge.get("color", "#FFFFFF"),
                    "font_size": badge.get("font_size", 11),
                    "font_weight": "bold",
                    "font_family": style.get("font_family", "Arial"),
                    "text_align": "center",
                    "vertical_align": "center",
                    "border": {"width": 0.1, "color": badge.get("background", "#005DAA")},
                    "padding": 0,
                    "auto_fit": False,
                },
            },
        )
    body_x = x + (badge_size + 34 if number else 72 if is_note_strip else 12)
    title, remaining_lines = _split_callout_title(body_lines)
    if is_note_strip:
        body_style_overrides = {"font_style": style.get("font_style", "italic") if style.get("italic", True) is not False else ""}
    else:
        body_style_overrides = {}
    body_style = {
        "color": style.get("body_color", style.get("color", "#1C1917")),
        "font_size": style.get("body_font_size", style.get("font_size", 11)),
        "font_weight": style.get("body_font_weight", "normal"),
        **body_style_overrides,
        "font_family": style.get("font_family", "Arial"),
        "text_align": "left",
        "vertical_align": style.get("vertical_align", "center" if is_note_strip else "top"),
        "padding": 0,
        "line_spacing": style.get("line_spacing", 0.92),
        "auto_fit": False,
    }
    title_style = {
        **body_style,
        "font_size": style.get("title_font_size", body_style["font_size"]),
        "font_weight": style.get("title_font_weight", "bold"),
        "vertical_align": "top",
    }
    if title and remaining_lines:
        runs = [
            {
                "text": f"{title} ",
                "font_weight": title_style.get("font_weight", "bold"),
                "font_size": title_style.get("font_size", body_style["font_size"]),
                "font_family": title_style.get("font_family", "Arial"),
                "color": title_style.get("color", "#1C1917"),
            },
            {
                "text": remaining_lines[0],
                "font_weight": body_style.get("font_weight", "normal"),
                "font_size": body_style.get("font_size", 11),
                "font_family": body_style.get("font_family", "Arial"),
                "color": body_style.get("color", "#1C1917"),
                "font_style": body_style.get("font_style", ""),
                "break_after": len(remaining_lines) > 1,
            },
        ]
        for line in remaining_lines[1:]:
            runs.append({
                "text": line,
                "font_weight": body_style.get("font_weight", "normal"),
                "font_size": body_style.get("font_size", 11),
                "font_family": body_style.get("font_family", "Arial"),
                "color": body_style.get("color", "#1C1917"),
                "font_style": body_style.get("font_style", ""),
                "break_after": line != remaining_lines[-1],
            })
        _add_textbox(
            slide,
            [body_x, y + 14, max(20, x + w - body_x - 14), max(20, h - 24)],
            [],
            body_style,
            runs,
        )
    else:
        _add_text(slide, [body_x, y + 14, max(20, x + w - body_x - 14), max(20, h - 24)], body_lines, body_style)


def _emit_chevron(slide, obj: dict) -> None:
    bbox = obj.get("bbox", [0, 0, 300, 120])
    style = obj.get("style", {})
    x, y, w, h = bbox
    _emit_shape(slide, {**obj, "text": [], "style": {**style, "shape_type": "chevron", "padding": 0}})
    lines = _get_text_lines(obj)
    if not lines:
        return
    text_bbox = [x + 28, y + h * 0.36, max(20, w - 66), h * 0.52]
    measurements = obj.get("reference_measurements")
    parts = measurements.get("component_parts") if isinstance(measurements, dict) else None
    if isinstance(parts, list):
        for part in parts:
            if not isinstance(part, dict) or part.get("type") != "label_text":
                continue
            part_bbox = part.get("bbox_px")
            if isinstance(part_bbox, list) and len(part_bbox) == 4:
                try:
                    text_bbox = [float(v) for v in part_bbox]
                except (TypeError, ValueError):
                    pass
            break
    text_style = {
        "color": style.get("color", "#FFFFFF"),
        "font_size": style.get("body_font_size", style.get("font_size", 12)),
        "font_weight": style.get("font_weight", "bold"),
        "font_family": style.get("font_family", "Arial"),
        "text_align": "center",
        "vertical_align": "center",
        "padding": 0,
        "line_spacing": style.get("line_spacing", 0.88),
        "auto_fit": False,
    }
    _add_text(slide, text_bbox, lines, text_style)


def _emit_nav(slide, obj: dict) -> None:
    bbox = obj.get("bbox", [0, 0, 1920, 44])
    style = obj.get("style", {})
    x, y, w, h = _px(bbox)
    bg = _style_color(style, "background", PALETTE["deep_blue"])
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _strip_theme_style(shp)
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    shp.line.fill.background()

    items = obj.get("items")
    if isinstance(items, list) and items:
        n = len(items)
        item_w = bbox[2] / n
        for i, item in enumerate(items):
            if not isinstance(item, dict):
                continue
            text = str(item.get("text", ""))
            if not text.strip():
                continue
            active = item.get("active", False)
            ix = bbox[0] + i * item_w
            iy = bbox[1]
            iw = item_w
            ih = bbox[3]
            item_style = {
                "color": style.get("color", "#FFFFFF"),
                "font_size": style.get("font_size", 11),
                "font_weight": "bold" if active else "normal",
                "text_align": "center",
                "vertical_align": "center",
            }
            _add_textbox(slide, [ix, iy, iw, ih], [text], item_style)
    else:
        text_lines = _get_text_lines(obj)
        if text_lines:
            nav_style = {
                "color": style.get("color", "#FFFFFF"),
                "font_size": style.get("font_size", 11),
                "text_align": "center",
                "vertical_align": "center",
            }
            _add_textbox(slide, bbox, text_lines, nav_style)


def _emit_takeaway(slide, obj: dict) -> None:
    bbox = obj.get("bbox", [0, 1000, 1920, 44])
    style = obj.get("style", {})
    x, y, w, h = _px(bbox)
    bg = _style_color(style, "background", PALETTE["deep_blue"])
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _strip_theme_style(shp)
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    shp.line.fill.background()

    label = obj.get("label", "")
    body = obj.get("body", "")
    if not label and not body:
        text_lines = _get_text_lines(obj)
        if text_lines:
            label_style = {
                "color": "#FFFFFF",
                "font_size": style.get("font_size", 11),
                "text_align": "left",
                "vertical_align": "center",
            }
            _add_textbox(slide, bbox, text_lines, label_style)
        return

    label_frac = style.get("label_width_fraction", 0.15)
    lw = bbox[2] * label_frac
    bw = bbox[2] * (1 - label_frac)

    if label:
        label_style = {
            "color": "#FFFFFF",
            "font_size": style.get("label_font_size", style.get("font_size", 11)),
            "font_weight": "bold",
            "text_align": "center",
            "vertical_align": "center",
        }
        _add_textbox(slide, [bbox[0], bbox[1], lw, bbox[3]], [label], label_style)

    div_x = bbox[0] + lw
    dx, dy, dw, dh = _px([div_x, bbox[1] + 4, 2, bbox[3] - 8])
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, dx, dy, dw, dh)
    _strip_theme_style(div)
    div.fill.solid()
    div.fill.fore_color.rgb = PALETTE["white"]
    div.line.fill.background()

    if body:
        body_style = {
            "color": "#FFFFFF",
            "font_size": style.get("body_font_size", style.get("font_size", 10)),
            "text_align": "left",
            "vertical_align": "center",
        }
        body_lines = body.split("\n") if isinstance(body, str) else [str(b) for b in body]
        _add_textbox(slide, [div_x + 10, bbox[1], bw - 10, bbox[3]], body_lines, body_style)


def _emit_table(slide, obj: dict) -> None:
    bbox = obj.get("bbox", [0, 0, 800, 400])
    style = obj.get("style", {})
    td = obj.get("table_data", {})
    headers = td.get("headers", [])
    rows = td.get("rows", [])
    if not rows and not headers:
        return

    all_rows = []
    if headers:
        all_rows.append(headers)
    all_rows.extend(rows)

    n_rows = len(all_rows)
    n_cols = max(len(r) for r in all_rows) if all_rows else 0
    if n_cols == 0:
        return

    x, y, w, h = _px(bbox)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table

    tblPr = tbl._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        sid = tblPr.find(qn("a:tableStyleId"))
        if sid is not None:
            tblPr.remove(sid)
        tblPr.set("bandRow", "0")
        tblPr.set("firstRow", "0")

    col_widths = td.get("col_widths")
    if col_widths and len(col_widths) >= n_cols:
        total_px = sum(col_widths[:n_cols]) or 1
        total_emu = int(bbox[2] * PX_TO_EMU_X)
        gridCols = tbl._tbl.findall(qn("a:tblGrid") + "/" + qn("a:gridCol"))
        for ci, gc in enumerate(gridCols):
            if ci < len(col_widths):
                gc.set("w", str(max(1, int(total_emu * col_widths[ci] / total_px))))

    row_heights = td.get("row_heights")
    if isinstance(row_heights, list) and len(row_heights) >= n_rows:
        total_px = sum(row_heights[:n_rows]) or 1
        total_emu = int(bbox[3] * PX_TO_EMU_Y)
        for ri, row in enumerate(tbl.rows):
            exact_height = max(1, int(total_emu * row_heights[ri] / total_px))
            row.height = Emu(exact_height)
            try:
                tbl._tbl.tr_lst[ri].set("h", str(exact_height))
            except Exception:
                pass

    header_bg = _parse_color(td.get("header_background", "#283E84")) or HEADER_BG
    header_color = _parse_color(td.get("header_color", "#FFFFFF")) or PALETTE["white"]
    row_header_bg = _parse_color(td.get("row_header_background", "#EAF2FB"))
    row_header_color = _parse_color(td.get("row_header_color", "#305496")) or PALETTE["deep_blue"]
    body_bg = _parse_color(td.get("body_background", style.get("background", "")))
    alt_bg = _parse_color(td.get("alternate_row_background", ""))
    border_color = _parse_color(td.get("border_color", "")) or _style_color(style, "border_color", BORDER_LIGHT) or BORDER_LIGHT
    border_width = int(max(BORDER_WIDTH_EMU, _float_value(td.get("border_width", style.get("line_width", 0.5)), 0.5) * 12700))
    header_border_width = int(max(border_width, _float_value(td.get("header_border_width", 1.0), 1.0) * 12700))
    cell_font_size = _float_value(td.get("font_size", style.get("font_size")), 10.0)
    header_font_size = _float_value(td.get("header_font_size", cell_font_size), cell_font_size)
    row_header_font_size = _float_value(td.get("row_header_font_size", cell_font_size), cell_font_size)
    font_family = str(style.get("font_family", td.get("font_family", "Arial")))
    cell_align = ALIGN_MAP.get(str(td.get("text_align", style.get("text_align", "left"))), PP_ALIGN.LEFT)
    header_align = ALIGN_MAP.get(str(td.get("header_text_align", td.get("text_align", "left"))), PP_ALIGN.LEFT)
    cell_padding = td.get("cell_padding", style.get("padding", 6))

    for ri, row_data in enumerate(all_rows):
        is_header = headers and ri == 0
        for ci in range(n_cols):
            cell_text = str(row_data[ci]) if ci < len(row_data) else ""
            cell = tbl.cell(ri, ci)
            is_row_header = ci == 0 and not is_header

            if is_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif is_row_header and row_header_bg:
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_header_bg
            elif alt_bg and ri % 2 == (0 if headers else 1):
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_bg
            elif body_bg:
                cell.fill.solid()
                cell.fill.fore_color.rgb = body_bg
            else:
                cell.fill.background()

            tcPr = cell._tc.find(qn("a:tcPr"))
            if tcPr is None:
                tcPr = SubElement(cell._tc, qn("a:tcPr"))
            _set_cell_border(tcPr, "a:lnT", border_color, border_width)
            _set_cell_border(tcPr, "a:lnB", border_color, header_border_width if is_header else border_width)
            _set_cell_border(tcPr, "a:lnL", border_color, border_width)
            _set_cell_border(tcPr, "a:lnR", border_color, border_width)
            v_align = str(td.get("vertical_align", style.get("vertical_align", "top")))
            tcPr.set("anchor", "ctr" if v_align == "center" else "b" if v_align == "bottom" else "t")

            tf = cell.text_frame
            _lock_text_body(
                tf,
                {
                    "vertical_align": "center" if v_align == "center" else "top",
                    "padding": cell_padding,
                    "auto_fit": style.get("auto_fit", td.get("auto_fit", False)),
                    "word_wrap": style.get("word_wrap", td.get("word_wrap", True)),
                },
                default_margin=6,
            )
            cell.text = ""

            p = tf.paragraphs[0]
            p.alignment = header_align if is_header else cell_align
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            run = p.add_run()
            run.text = cell_text
            run.font.name = font_family
            _lock_run_typeface(run, font_family)
            run.font.size = Pt(header_font_size if is_header else row_header_font_size if is_row_header else cell_font_size)
            if is_header:
                run.font.color.rgb = header_color
                run.font.bold = True
            elif is_row_header:
                run.font.color.rgb = row_header_color
                run.font.bold = True
            else:
                run.font.color.rgb = _style_color(style, "color", PALETTE["text_dark"]) or PALETTE["text_dark"]
                run.font.bold = str(style.get("font_weight", "normal")) in ("bold", "700", "800", "900")


def _emit_chart(slide, obj: dict, errors: list[str]) -> None:
    bbox = obj.get("bbox", [0, 0, 800, 400])
    cd = obj.get("chart_data", {})
    ct_str = str(cd.get("chart_type", "")).lower().strip()
    if ct_str in CHART_ALIASES:
        ct_str = CHART_ALIASES[ct_str]
    if ct_str not in CHART_TYPE_MAP:
        errors.append(f"[CHART] unsupported type '{ct_str}' in object {obj.get('id', '?')}")
        return

    cats = cd.get("categories", [])
    series_list = cd.get("series", [])
    if not cats or not series_list:
        errors.append(f"[CHART] missing categories or series in object {obj.get('id', '?')}")
        return

    if len(series_list) > MAX_SERIES:
        series_list = series_list[:MAX_SERIES]

    valid_series = []
    for s in series_list:
        if not isinstance(s, dict):
            continue
        vals = s.get("values")
        if not isinstance(vals, list):
            continue
        if ct_str != "scatter" and len(vals) != len(cats):
            if len(vals) < len(cats):
                vals = vals + [None] * (len(cats) - len(vals))
            else:
                vals = vals[:len(cats)]
            s = {**s, "values": vals}
        valid_series.append(s)

    if not valid_series:
        errors.append(f"[CHART] no valid series in object {obj.get('id', '?')}")
        return

    if ct_str == "scatter":
        chart_data = XyChartData()
        for s in valid_series:
            xy = chart_data.add_series(s.get("name", ""))
            for pt in s.get("values", []):
                if isinstance(pt, list) and len(pt) >= 2:
                    xy.add_data_point(pt[0], pt[1])
    else:
        chart_data = CategoryChartData()
        chart_data.categories = [str(c) for c in cats]
        for s in valid_series:
            chart_data.add_series(s.get("name", ""), s.get("values", []))

    x, y, w, h = _px(bbox)
    try:
        chart_shape = slide.shapes.add_chart(CHART_TYPE_MAP[ct_str], x, y, w, h, chart_data)
    except Exception as exc:
        errors.append(f"[CHART] add_chart failed: {exc}")
        return

    chart = chart_shape.chart
    title = str(cd.get("title", "")).strip()
    title_font_size = _float_value(cd.get("title_font_size"), 10.0)
    if title:
        try:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].runs[0].font.name = "Arial"
            chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(title_font_size)
            chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
        except Exception:
            pass
    _set_chart_manual_layout(chart, cd.get("plot_area", {"x": 0.08, "y": 0.15, "w": 0.84, "h": 0.70}))
    if ct_str in ("column", "bar", "column_stacked", "bar_stacked"):
        _set_chart_bar_geometry(chart, cd.get("bar_gap_width", 80), cd.get("bar_overlap"))
    try:
        for i, ps in enumerate(chart.series):
            if i >= MAX_SERIES:
                break
            ps.format.fill.solid()
            ps.format.fill.fore_color.rgb = SERIES_COLORS[i]
            if ct_str == "line":
                ps.format.line.color.rgb = SERIES_COLORS[i]
                ps.format.line.width = Emu(28575)
    except Exception:
        pass

    try:
        chart.font.name = "Arial"
        chart.font.size = Pt(_float_value(cd.get("font_size"), 8.5))
    except Exception:
        pass

    try:
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.name = "Arial"
        cat_axis.tick_labels.font.size = Pt(_float_value(cd.get("category_font_size"), 7.5))
        cat_axis.has_major_gridlines = False
    except Exception:
        pass

    try:
        val_axis = chart.value_axis
        val_axis.tick_labels.font.name = "Arial"
        val_axis.tick_labels.font.size = Pt(_float_value(cd.get("value_font_size"), 7.5))
        val_axis.has_major_gridlines = bool(cd.get("major_gridlines", True))
        val_axis.major_gridlines.format.line.color.rgb = BORDER_LIGHT
        val_axis.major_gridlines.format.line.width = Emu(BORDER_WIDTH_EMU)
    except Exception:
        pass

    show_legend = len(valid_series) >= 2
    chart.has_legend = show_legend
    if show_legend:
        try:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.name = "Arial"
            chart.legend.font.size = Pt(9)
        except Exception:
            pass

    if cd.get("data_labels") and ct_str in ("column", "bar", "line", "column_stacked", "bar_stacked", "area"):
        try:
            chart.plots[0].has_data_labels = True
        except Exception:
            pass
        for ps in chart.series:
            try:
                dl = ps.data_labels
                dl.show_value = True
                dl.font.name = "Arial"
                dl.font.size = Pt(_float_value(cd.get("data_label_font_size"), 8.0))
                dl.font.bold = True
            except Exception:
                pass

    if ct_str in ("pie", "doughnut"):
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.show_category_name = True
            dl.show_percentage = True
            dl.font.name = "Arial"
            dl.font.size = Pt(10)
        except Exception:
            pass


def _emit_divider(slide, obj: dict) -> None:
    bbox = obj.get("bbox", [0, 0, 100, 2])
    style = obj.get("style", {})
    x, y, w, h = _px(bbox)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _strip_theme_style(shp)
    fill = _style_color(style, "background", BORDER_LIGHT)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()


def _set_connector_arrow(shp, style: dict) -> None:
    end_arrow = str(style.get("end_arrow") or style.get("arrow") or "").strip().lower()
    begin_arrow = str(style.get("begin_arrow") or "").strip().lower()
    if end_arrow not in {"triangle", "arrow", "end", "true"} and begin_arrow not in {"triangle", "arrow", "begin", "true"}:
        return
    try:
        sp_pr = shp._element.find(qn("p:spPr"))
        ln = sp_pr.find(qn("a:ln")) if sp_pr is not None else None
        if ln is None and sp_pr is not None:
            ln = SubElement(sp_pr, qn("a:ln"))
        if ln is None:
            return
        for child in list(ln):
            if child.tag in {qn("a:headEnd"), qn("a:tailEnd")}:
                ln.remove(child)
        if begin_arrow in {"triangle", "arrow", "begin", "true"}:
            head = SubElement(ln, qn("a:headEnd"))
            head.set("type", "triangle")
        if end_arrow in {"triangle", "arrow", "end", "true"}:
            tail = SubElement(ln, qn("a:tailEnd"))
            tail.set("type", "triangle")
    except Exception:
        pass


def _connector_points(obj: dict) -> list[tuple[float, float]]:
    points = obj.get("points")
    if isinstance(points, list):
        parsed: list[tuple[float, float]] = []
        for point in points:
            if (
                isinstance(point, list)
                and len(point) >= 2
                and isinstance(point[0], (int, float))
                and isinstance(point[1], (int, float))
            ):
                parsed.append((float(point[0]), float(point[1])))
        if len(parsed) >= 2:
            return parsed
    bbox = obj.get("bbox", [0, 0, 100, 0])
    if isinstance(bbox, list) and len(bbox) >= 4:
        x, y, w, h = bbox[:4]
        if all(isinstance(v, (int, float)) for v in (x, y, w, h)):
            return [(float(x), float(y)), (float(x + w), float(y + h))]
    return [(0.0, 0.0), (100.0, 0.0)]


def _emit_connector(slide, obj: dict) -> None:
    style = obj.get("style", {})
    points = _connector_points(obj)
    color = _style_color(style, "color", _style_color(style, "border_color", PALETTE["deep_blue"]))
    width_px = style.get("width", style.get("line_width", 2))
    conn_type = MSO_CONNECTOR.CURVE if str(style.get("connector_type", "")).lower() == "curve" else MSO_CONNECTOR.STRAIGHT
    for idx in range(len(points) - 1):
        x1, y1 = points[idx]
        x2, y2 = points[idx + 1]
        shp = slide.shapes.add_connector(
            conn_type,
            Emu(int(x1 * PX_TO_EMU_X)),
            Emu(int(y1 * PX_TO_EMU_Y)),
            Emu(int(x2 * PX_TO_EMU_X)),
            Emu(int(y2 * PX_TO_EMU_Y)),
        )
        _apply_line_style(
            shp.line,
            {**style, "line_color": style.get("color", style.get("border_color")), "line_width": width_px},
            default_color=color,
            default_width_px=_float_value(width_px, 2.0),
        )
        if idx == len(points) - 2:
            _set_connector_arrow(shp, style)

    label = _get_text_lines(obj)
    if label:
        _add_textbox(slide, obj.get("bbox", [0, 0, 100, 30]), label, style, _get_text_runs(obj))


def _emit_accent_border(slide, bbox: list, accent: dict) -> None:
    """Emit a thin colored rectangle on one side of a card (left/right/top/bottom)."""
    color = _parse_color(accent.get("color", ""))
    if not color:
        return
    accent_w_px = _float_value(accent.get("width"), 4)
    side = str(accent.get("side", "left")).lower()
    bx, by, bw, bh = bbox

    if side == "left":
        ax, ay, aw, ah = bx, by, accent_w_px, bh
    elif side == "right":
        ax, ay, aw, ah = bx + bw - accent_w_px, by, accent_w_px, bh
    elif side == "top":
        ax, ay, aw, ah = bx, by, bw, accent_w_px
    elif side == "bottom":
        ax, ay, aw, ah = bx, by + bh - accent_w_px, bw, accent_w_px
    else:
        ax, ay, aw, ah = bx, by, accent_w_px, bh

    sx, sy, sw, sh = _px([ax, ay, aw, ah])
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx, sy, sw, sh)
    _strip_theme_style(shp)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False


def _emit_image(slide, obj: dict) -> None:
    """Embed an external image file (icon, small asset) at the given bbox."""
    bbox = obj.get("bbox", [0, 0, 100, 100])
    img_path = obj.get("path") or obj.get("style", {}).get("path", "")
    if not img_path:
        return
    p = Path(img_path)
    if not p.exists():
        return
    img_w, img_h = bbox[2], bbox[3]
    if (img_w * img_h) / (HTML_W * HTML_H) > FULL_SLIDE_IMAGE_AREA_RATIO:
        return
    x, y, w, h = _px(bbox)
    slide.shapes.add_picture(str(p), x, y, width=w, height=h)


def _emit_icon(slide, obj: dict) -> None:
    bbox = obj.get("bbox", [0, 0, 48, 48])
    style = obj.get("style", {})
    img_path = obj.get("path") or style.get("path", "")
    if img_path:
        _emit_image(slide, obj)
        return
    shape_name = str(style.get("shape_type", "oval")).strip().lower()
    icon_obj = {**obj, "style": {**style, "shape_type": shape_name}}
    if shape_name not in {"oval", "chevron"}:
        icon_obj["style"]["shape_type"] = "shape"
    _emit_shape(slide, icon_obj)


def _compile_slide(prs: Presentation, graph: dict, errors: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    bg_color = _parse_color(graph.get("background"))
    if bg_color and bg_color != PALETTE["white"]:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

    objects = graph.get("objects", [])
    if isinstance(objects, list):
        objects = sorted(
            objects,
            key=lambda obj: obj.get("z_order", 0) if isinstance(obj, dict) else 0,
        )
    for obj in objects:
        kind = str(obj.get("object_kind", "")).strip()
        if kind == "kpi":
            _emit_kpi(slide, obj)
        elif kind == "callout":
            _emit_callout(slide, obj)
        elif kind == "chevron":
            _emit_chevron(slide, obj)
        elif kind in ("shape", "shape_group", "card", "panel", "badge", "nav_indicator", "callout_strip"):
            _emit_shape(slide, obj)
        elif kind in ("text", "text_box", "title", "source"):
            _emit_text(slide, obj)
        elif kind == "nav":
            _emit_nav(slide, obj)
        elif kind == "takeaway":
            _emit_takeaway(slide, obj)
        elif kind in ("table", "native_table"):
            _emit_table(slide, obj)
        elif kind in ("chart", "native_chart"):
            _emit_chart(slide, obj, errors)
        elif kind == "divider":
            _emit_divider(slide, obj)
        elif kind == "connector":
            _emit_connector(slide, obj)
        elif kind == "image":
            _emit_image(slide, obj)
        elif kind == "icon":
            _emit_icon(slide, obj)
        else:
            _emit_shape(slide, obj)


def compile_deck(graph_paths: list, pptx_path) -> dict:
    """Compile object graph JSON files into a single PPTX.

    Returns summary dict with per-slide info and any errors.
    """
    graph_paths = [Path(p) for p in graph_paths]
    pptx_path = Path(pptx_path)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    summary: dict = {"slides": [], "errors": [], "pptx": str(pptx_path)}

    for idx, gpath in enumerate(graph_paths, 1):
        graph = json.loads(gpath.read_text(encoding="utf-8"))
        slide_errors: list[str] = []
        objects = graph.get("objects", [])
        summary["slides"].append({
            "idx": idx,
            "graph": str(gpath),
            "objects": len(objects),
        })
        _compile_slide(prs, graph, slide_errors)
        if slide_errors:
            summary["errors"].append({"idx": idx, "items": slide_errors})

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    return summary
