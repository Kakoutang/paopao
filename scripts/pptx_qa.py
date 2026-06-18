"""
pptx_qa.py — Post-render PPTX inspection.

Catches mechanical bugs that escape composer + render:
  - Shape with opaque fill covering chart/table
  - Shape extending past canvas (1920×1080 / 13.33"×7.5")
  - Shape extending into the detected bottom takeaway band
  - Chart with 0 dimension or invalid data

For each issue: report it. Strict delivery must refill HTML/assets and rerender
instead of mutating the finished PPTX as a separate repair layer.

Usage:
    from pptx_qa import qa_and_fix
    summary = qa_and_fix(pptx_path)  # inspect only by default

Pipeline integration: pipeline_v3 stage 4.7 (after render_pptx, before download).
"""
from __future__ import annotations

import json
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

# Canvas constants (16:9 PPTX defaults)
CANVAS_W_EMU = 12192000   # 13.333 inches
CANVAS_H_EMU = 6858000    # 7.5 inches
FALLBACK_TAKEAWAY_TOP_EMU = int(CANVAS_H_EMU * 0.86)
TAKEAWAY_MIN_H_EMU = 355600  # ~56px at 96 DPI
FULL_SLIDE_PICTURE_AREA_RATIO = 0.82
FULL_SLIDE_PICTURE_EDGE_TOL_EMU = 304800  # ~0.33in, maps to ~48px on 1920 canvas

# Tolerance: shapes overlap if their bbox intersects, regardless of fill type.
# But only treat as "covering" issue if the offending shape has an opaque fill.


@dataclass
class QAIssue:
    slide_idx: int
    kind: str           # 'cover_chart', 'cover_table', 'overflow_canvas', 'overflow_takeaway', 'chart_zero_dim', 'full_slide_picture'
    detail: str
    fixed: bool = False
    fix_action: str = ""


@dataclass
class QASummary:
    pptx_path: str
    issues: list[QAIssue] = field(default_factory=list)
    fixes_applied: int = 0

    def to_dict(self) -> dict[str, Any]:
        return {
            "pptx_path": self.pptx_path,
            "total_issues": len(self.issues),
            "fixes_applied": self.fixes_applied,
            "issues": [
                {
                    "slide": i.slide_idx + 1,
                    "kind": i.kind,
                    "detail": i.detail,
                    "fixed": i.fixed,
                    "fix_action": i.fix_action,
                }
                for i in self.issues
            ],
        }


def _bboxes_intersect(a, b) -> bool:
    """a, b: (x, y, w, h) tuples in EMU."""
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    if ax + aw <= bx or bx + bw <= ax: return False
    if ay + ah <= by or by + bh <= ay: return False
    return True


def _shape_bbox(sh) -> tuple[int, int, int, int] | None:
    if sh.left is None or sh.top is None or sh.width is None or sh.height is None:
        return None
    return (sh.left, sh.top, sh.width, sh.height)


def _is_opaque_fill(sh) -> bool:
    """Check if shape's fill is solid (covers content below). Background = transparent."""
    try:
        ft = sh.fill.type
        # 1=SOLID, 2=PATTERNED, 3=GRADIENT, 4=PICTURE, 5=BACKGROUND (transparent), 6=GROUP
        return ft in (1, 2, 3, 4)
    except Exception:
        return False


def _solid_fill_hex(sh) -> str:
    try:
        if sh.fill.type != 1:
            return ""
        rgb = sh.fill.fore_color.rgb
        return str(rgb).upper() if rgb is not None else ""
    except Exception:
        return ""


def _detect_takeaway_top(slide) -> tuple[int, bool]:
    """Return the top of the actual bottom takeaway band.

    The system style wants a deep-blue bottom strip that can be taller than
    older QA allowed. Detect it from slide geometry/color instead of forcing
    all decks into a fixed 1014px top line.
    """
    candidates: list[tuple[int, int, int, int, int]] = []
    for sh in slide.shapes:
        bb = _shape_bbox(sh)
        if not bb:
            continue
        x, y, w, h = bb
        fill_hex = _solid_fill_hex(sh)
        if (
            fill_hex == "305496"
            and y >= int(CANVAS_H_EMU * 0.72)
            and w >= int(CANVAS_W_EMU * 0.45)
            and h >= int(TAKEAWAY_MIN_H_EMU * 0.55)
        ):
            candidates.append((y, x, w, h, y + h))
    if not candidates:
        return FALLBACK_TAKEAWAY_TOP_EMU, False
    candidates.sort()
    return candidates[0][0], True


def _make_transparent(sh) -> bool:
    """Set shape fill to transparent. Returns True if change applied."""
    try:
        sh.fill.background()
        return True
    except Exception:
        return False


def _is_picture_shape(sh) -> bool:
    try:
        return sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return False


def _is_full_slide_picture_bbox(bb: tuple[int, int, int, int]) -> bool:
    x, y, w, h = bb
    area_ratio = (w * h) / float(CANVAS_W_EMU * CANVAS_H_EMU)
    near_edges = (
        x <= FULL_SLIDE_PICTURE_EDGE_TOL_EMU
        and y <= FULL_SLIDE_PICTURE_EDGE_TOL_EMU
        and x + w >= CANVAS_W_EMU - FULL_SLIDE_PICTURE_EDGE_TOL_EMU
        and y + h >= CANVAS_H_EMU - FULL_SLIDE_PICTURE_EDGE_TOL_EMU
    )
    return area_ratio >= FULL_SLIDE_PICTURE_AREA_RATIO and near_edges


def _clip_bbox_to_canvas(sh) -> tuple[bool, str]:
    """If shape extends past canvas edges, shrink it to fit. Returns (changed, reason)."""
    bb = _shape_bbox(sh)
    if not bb: return False, ""
    x, y, w, h = bb
    new_x, new_y, new_w, new_h = x, y, w, h
    actions = []
    if x < 0:
        new_w = max(1, w + x)
        new_x = 0
        actions.append(f"x {x}→0")
    if y < 0:
        new_h = max(1, h + y)
        new_y = 0
        actions.append(f"y {y}→0")
    if new_x + new_w > CANVAS_W_EMU:
        new_w = max(1, CANVAS_W_EMU - new_x)
        actions.append(f"w clipped to {new_w}")
    if new_y + new_h > CANVAS_H_EMU:
        new_h = max(1, CANVAS_H_EMU - new_y)
        actions.append(f"h clipped to {new_h}")
    if not actions:
        return False, ""
    sh.left = Emu(new_x); sh.top = Emu(new_y); sh.width = Emu(new_w); sh.height = Emu(new_h)
    return True, ", ".join(actions)


def _clip_bbox_above_takeaway(sh, takeaway_top: int) -> tuple[bool, str]:
    """If shape's bottom extends past takeaway top, shrink h."""
    bb = _shape_bbox(sh)
    if not bb: return False, ""
    x, y, w, h = bb
    if y + h <= takeaway_top:
        return False, ""
    if y >= takeaway_top:
        # Shape entirely below takeaway — leave it alone (probably the takeaway itself or source)
        return False, ""
    new_h = max(1, takeaway_top - y)
    sh.height = Emu(new_h)
    return True, f"h clipped from {h} to {new_h} (was extending into detected takeaway band)"


def qa_and_fix(pptx_path: str | Path, *, verbose: bool = False, fix: bool = False) -> QASummary:
    pptx_path = Path(pptx_path)
    p = Presentation(pptx_path)
    summary = QASummary(pptx_path=str(pptx_path))

    for slide_idx, slide in enumerate(p.slides):
        # 0. Formal Spark path allows localized preserved assets, but never a
        # full-slide picture layer. This catches decks generated outside the
        # HTML renderer too.
        for j, sh in enumerate(slide.shapes):
            if not _is_picture_shape(sh):
                continue
            bb = _shape_bbox(sh)
            if not bb:
                continue
            if _is_full_slide_picture_bbox(bb):
                summary.issues.append(QAIssue(
                    slide_idx=slide_idx,
                    kind="full_slide_picture",
                    detail=f"picture[{j}] bbox={bb} appears to be a full-slide image layer",
                    fixed=False,
                    fix_action="blocked_for_formal_delivery",
                ))

        # 1. Find charts and tables (the "underlay" content that mustn't be covered)
        underlay = []  # (kind, idx, bbox)
        for i, sh in enumerate(slide.shapes):
            bb = _shape_bbox(sh)
            if not bb: continue
            if sh.has_chart:
                underlay.append(("chart", i, bb))
            elif sh.has_table:
                underlay.append(("table", i, bb))

        # 2. For each underlay, find shapes ABOVE it (later in z-order) that are opaque + overlap
        for kind, u_idx, u_bb in underlay:
            for j, sh in enumerate(slide.shapes):
                if j <= u_idx: continue  # below in z-order, can't cover
                if sh.has_chart or sh.has_table: continue
                bb = _shape_bbox(sh)
                if not bb: continue
                if not _bboxes_intersect(u_bb, bb): continue
                if not _is_opaque_fill(sh): continue
                # Skip small badge/pill shapes (score indicators inside tables).
                # These are intentionally opaque overlays on table cells.
                # Threshold: shapes smaller than ~1 inch in both dimensions.
                _, _, sw, sh_h = bb
                if sw < 914400 and sh_h < 914400:  # 914400 EMU = 1 inch
                    continue
                detail = f"shape[{j}] (bbox {bb}) opaque-fills over {kind}[{u_idx}]"
                fixed = _make_transparent(sh) if fix else False
                summary.issues.append(QAIssue(
                    slide_idx=slide_idx,
                    kind=f"cover_{kind}",
                    detail=detail,
                    fixed=fixed,
                    fix_action="set_fill_transparent" if fixed else "refill_html_required",
                ))
                if fixed: summary.fixes_applied += 1

        # 3. Find shapes overflowing canvas
        for j, sh in enumerate(slide.shapes):
            bb = _shape_bbox(sh)
            if not bb: continue
            x, y, w, h = bb
            if x < 0 or y < 0 or x + w > CANVAS_W_EMU or y + h > CANVAS_H_EMU:
                changed, reason = _clip_bbox_to_canvas(sh) if fix else (False, "refill_html_required")
                summary.issues.append(QAIssue(
                    slide_idx=slide_idx,
                    kind="overflow_canvas",
                    detail=f"shape[{j}] bbox=({x},{y},{w},{h}) outside canvas",
                    fixed=changed,
                    fix_action=f"clip: {reason}" if changed else reason,
                ))
                if changed: summary.fixes_applied += 1

        # 4. Find shapes overflowing the takeaway region. The top is detected
        # from the actual deep-blue bottom band instead of a fixed y value.
        takeaway_top, takeaway_detected = _detect_takeaway_top(slide)
        for j, sh in enumerate(slide.shapes):
            bb = _shape_bbox(sh)
            if not bb: continue
            x, y, w, h = bb
            if y < takeaway_top and y + h > takeaway_top:
                # Don't fix if shape is the title band, nav, or covers full canvas (pinning element)
                if w >= CANVAS_W_EMU * 0.95 and h >= CANVAS_H_EMU * 0.7:
                    continue  # whole-slide background, leave alone
                changed, reason = _clip_bbox_above_takeaway(sh, takeaway_top) if fix else (False, "refill_html_required")
                summary.issues.append(QAIssue(
                    slide_idx=slide_idx,
                    kind="overflow_takeaway",
                    detail=(
                        f"shape[{j}] bbox=({x},{y},{w},{h}) extends into "
                        f"{'detected' if takeaway_detected else 'fallback'} takeaway band top={takeaway_top}"
                    ),
                    fixed=changed,
                    fix_action=f"clip: {reason}" if changed else reason,
                ))
                if changed: summary.fixes_applied += 1

        # 5. Find charts with 0 dimension
        for j, sh in enumerate(slide.shapes):
            if not sh.has_chart: continue
            bb = _shape_bbox(sh)
            if not bb: continue
            x, y, w, h = bb
            if w < 50 * 9525 or h < 50 * 9525:  # 50px in EMU (9525 EMU/px @ 96 DPI)
                summary.issues.append(QAIssue(
                    slide_idx=slide_idx,
                    kind="chart_zero_dim",
                    detail=f"chart[{j}] has tiny bbox ({w}, {h}) — likely render error",
                    fixed=False,
                    fix_action="manual_review_needed",
                ))

    if summary.fixes_applied > 0:
        p.save(pptx_path)
        if verbose:
            print(f"  pptx_qa: {summary.fixes_applied} auto-fixes applied to {pptx_path}")

    if verbose:
        for issue in summary.issues:
            mark = "✅" if issue.fixed else "⚠"
            print(f"    {mark} slide{issue.slide_idx+1} {issue.kind}: {issue.detail[:80]}")

    return summary


# ── CLI ───────────────────────────────────────────────────────────
if __name__ == "__main__":
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx_path")
    ap.add_argument("--verbose", "-v", action="store_true")
    ap.add_argument("--fix", action="store_true", help="Mutate PPTX mechanically; not allowed for final delivery")
    args = ap.parse_args()
    summary = qa_and_fix(args.pptx_path, verbose=args.verbose, fix=args.fix)
    print(json.dumps(summary.to_dict(), indent=2, ensure_ascii=False))
